import os
import pathlib
import re
from pptx import Presentation
from docx import Document

RAW_DATA_DIR = pathlib.Path("./data/EXE101")
PROCESSED_DATA_DIR = pathlib.Path("./data/processed")

class MarkdownParser:
    def __init__(self):
        self.list_counter = 0

    def clean_and_format(self, text):
        if not text:
            return ""
        
        # MD010: Fix Hard Tabs
        text = text.replace('\t', '    ')
        
        # MD034: Fix Bare URLs
        url_pattern = r'(https?://[^\s<>"]+|www\.[^\s<>"]+)'
        text = re.sub(url_pattern, r'<\1>', text)
        
        # MD009: Đảm bảo 2 spaces cho line break, tránh lỗi "Actual: 1"
        return text.strip() + "  "

    def handle_ordered_list(self, text):
        """
        Phát hiện và đánh số lại danh sách tăng dần (MD029)
        """
        # Kiểm tra xem dòng có phải là danh sách số không (ví dụ: 1. hoặc 1))
        match = re.match(r'^\s*\d+[.)]\s+(.*)', text)
        
        if match:
            self.list_counter += 1
            content = match.group(1)
            # Trả về số thứ tự đã được format tăng dần
            return f"{self.list_counter}. {content}"
        else:
            # Nếu không phải danh sách, reset bộ đếm
            self.list_counter = 0
            return text

    def parse_pptx(self, file_path):
        try:
            prs = Presentation(file_path)
            md_output = [f"# {file_path.stem}  "]
            
            for i, slide in enumerate(prs.slides):
                md_output.append(f"\n## Slide {i+1}  \n")
                self.list_counter = 0 # Reset mỗi slide
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        cleaned = self.clean_and_format(shape.text)
                        formatted = self.handle_ordered_list(cleaned)
                        
                        # MD032: Blank lines around lists
                        if formatted.startswith(tuple("123456789")) or formatted.startswith(('*', '-', '+')):
                            md_output.append(f"\n{formatted}\n")
                        else:
                            md_output.append(formatted)
                
                md_output.append("\n\n---\n")
            return "\n".join(md_output)
        except Exception as e:
            return f"Error parsing PPTX {file_path.name}: {e}"

    def parse_docx(self, file_path):
        try:
            doc = Document(file_path)
            md_output = [f"# {file_path.stem}  \n"]
            self.list_counter = 0 
            
            for para in doc.paragraphs:
                if para.text.strip():
                    cleaned = self.clean_and_format(para.text)
                    formatted = self.handle_ordered_list(cleaned)
                    
                    if formatted.lstrip().startswith(tuple("123456789")) or formatted.lstrip().startswith(('*', '-', '+')):
                        md_output.append(f"\n{formatted}\n")
                    else:
                        md_output.append(formatted + "\n")
                        
            return "\n".join(md_output)
        except Exception as e:
            return f"Error parsing DOCX {file_path.name}: {e}"

def main():
    parser = MarkdownParser()
    PROCESSED_DATA_DIR.mkdir(parents=True, exist_ok=True)
    file_count = 0
    
    for root, _, files in os.walk(RAW_DATA_DIR):
        for file in files:
            file_path = pathlib.Path(root) / file
            suffix = file_path.suffix.lower()
            
            if suffix == ".pptx":
                content = parser.parse_pptx(file_path)
            elif suffix == ".docx":
                content = parser.parse_docx(file_path)
            else:
                continue

            output_file = PROCESSED_DATA_DIR / f"{file_path.stem}.md"
            with open(output_file, "w", encoding="utf-8") as f:
                # Xử lý MD010 & MD032 cuối cùng
                content = content.replace('\t', '    ')
                sanitized = re.sub(r'\n{3,}', '\n\n', content)
                
                f.write(sanitized.strip())
                f.write("\n") # MD047
            file_count += 1

    print(f"Done! Processed {file_count} files.")

if __name__ == "__main__":
    main()