import os
import pathlib
import re
from pptx import Presentation
from docx import Document

RAW_DATA_DIR = pathlib.Path("./data/EXE101")
PROCESSED_DATA_DIR = pathlib.Path("./data/processed")

def clean_text(text):
    if not text:
        return ""
    
    # MD034: Fix bare URLs by wrapping them in angle brackets
    url_pattern = r'(https?://[^\s<>"]+|www\.[^\s<>"]+)'
    text = re.sub(url_pattern, r'<\1>', text)
    
    # MD009: Trailing spaces
    return text.strip()

def pptx_parser(file_path):
    try:
        prs = Presentation(file_path)
        md_output = [f"# {file_path.stem}  "] # MD009
        
        for i, slide in enumerate(prs.slides):
            md_output.append(f"\n## Slide {i+1}  \n")
            
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    cleaned = clean_text(shape.text)
                    
                    # MD032: List not surrounded by blank space
                    if cleaned.lstrip().startswith(('*', '-', '+', '1.')):
                        slide_text.append(f"\n{cleaned}\n")
                    else:
                        slide_text.append(cleaned)
            
            if slide_text:
                md_output.append("\n\n".join(slide_text))
            
            md_output.append("\n\n---\n")
            
        return "\n".join(md_output)
    except Exception as e:
        return f"Error parsing PPTX {file_path.name}: {e}"

def docx_parser(file_path):
    try:
        doc = Document(file_path)
        md_output = [f"# {file_path.stem}  \n"]
        
        for para in doc.paragraphs:
            text = clean_text(para.text)
            if text.strip():
                # MD032
                if text.lstrip().startswith(('*', '-', '+', '1.')):
                    md_output.append(f"\n{text}\n")
                else:
                    md_output.append(text + "\n")
                
        return "\n".join(md_output)
    except Exception as e:
        return f"Error parsing DOCX {file_path.name}: {e}"

def main():
    PROCESSED_DATA_DIR.mkdir(parents=True, exist_ok=True)
    file_count = 0
    
    for root, _, files in os.walk(RAW_DATA_DIR):
        for file in files:
            file_path = pathlib.Path(root) / file
            suffix = file_path.suffix.lower()
            
            content = ""
            if suffix == ".pptx":
                print(f"Converting PPTX: {file}")
                content = pptx_parser(file_path)
            elif suffix == ".docx":
                print(f"Converting DOCX: {file}")
                content = docx_parser(file_path)
            else:
                continue

            output_file = PROCESSED_DATA_DIR / f"{file_path.stem}.md"
            with open(output_file, "w", encoding="utf-8") as f:
                sanitized_content = re.sub(r'\n{3,}', '\n\n', content)
                f.write(sanitized_content.strip())
                # MD049: No endline characters at the end of file
                f.write("\n")
            file_count += 1

    print(f"\nExtracted {file_count} files to {PROCESSED_DATA_DIR}")

if __name__ == "__main__":
    main()